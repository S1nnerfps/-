"""
Personal CS Coach 答辩 PPT 终稿
核心卖点：预调制 CS2 战术 AI Agent
用户购买 Agent → 自行接入 DeepSeek API
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# === 配色 ===
BG      = RGBColor(0x0D, 0x11, 0x17)
SURFACE = RGBColor(0x16, 0x1B, 0x22)
CYAN    = RGBColor(0x58, 0xA6, 0xFF)
YELLOW  = RGBColor(0xF5, 0xC8, 0x42)
WHITE   = RGBColor(0xE6, 0xED, 0xF3)
DIM     = RGBColor(0x8B, 0x94, 0x9E)
ORANGE  = RGBColor(0xE8, 0x75, 0x1A)
GREEN   = RGBColor(0x3A, 0x8C, 0x3C)
PURPLE  = RGBColor(0x9B, 0x3B, 0xC4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

def bg(slide): slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG
def rect(slide, l, t, w, h, color): s = slide.shapes.add_shape(1, l, t, w, h); s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background(); return s
def txt(slide, l, t, w, h, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    b = slide.shapes.add_textbox(l, t, w, h); b.text_frame.word_wrap = True
    p = b.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold; p.alignment = align; return b
def bullets(slide, l, t, w, h, items, size=13, color=WHITE, bullet='•'):
    b = slide.shapes.add_textbox(l, t, w, h); b.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = b.text_frame.paragraphs[0] if i == 0 else b.text_frame.add_paragraph()
        p.text = f"{bullet} {item}"; p.font.size = Pt(size); p.font.color.rgb = color; p.space_after = Pt(5)
    return b
def bar(slide, l, t, h, color=CYAN, w=Inches(0.06)): return rect(slide, l, t, w, h, color)
def hr(slide, l, t, w, color=CYAN, h=Inches(0.025)): return rect(slide, l, t, w, h, color)
def pg(slide, n): txt(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35), str(n), size=10, color=DIM, align=PP_ALIGN.RIGHT)
def footer(slide, text="Personal CS Coach — Agent + API 架构  |  CS2 战术智能"): txt(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3), text, size=9, color=DIM)

# ============================
# 0. 封面
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.05), CYAN)
txt(s, Inches(1.0), Inches(1.2), Inches(11), Inches(1.3), "Personal CS Coach", size=72, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(2.6), Inches(11), Inches(0.7), "预调制 CS2 战术 AI · Agent 驱动", size=24, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
hr(s, Inches(3.5), Inches(3.5), Inches(6.3), CYAN)
txt(s, Inches(1.5), Inches(3.8), Inches(10), Inches(0.5), "预训练战术模型  +  用户自接 DeepSeek API  +  本地推理闭环", size=17, color=DIM, align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(4.5), Inches(10), Inches(0.4), "设计思维期末大作业  ·  产品路演", size=17, color=DIM, align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4), "2026 年 6 月 23 日", size=14, color=DIM, align=PP_ALIGN.CENTER)
pg(s, 0)

# ============================
# 1. 目录
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.7), Inches(1.2), Inches(0.5))
txt(s, Inches(1.0), Inches(1.25), Inches(5), Inches(0.5), "CONTENTS", size=14, color=CYAN, bold=True)
txt(s, Inches(1.0), Inches(1.8), Inches(10), Inches(0.7), "答辩内容概览", size=34, color=WHITE, bold=True)
items = [
    ("01", "为什么要做", "项目背景  /  用户需求  /  核心痛点", CYAN),
    ("02", "我们做了什么", "Agent 架构  /  产品方案  /  Demo 演示", ORANGE),
    ("03", "商业可行性", "Agent 卖点  /  市场与竞品  /  商业模式", GREEN),
    ("04", "团队与展望", "预训练路线图  /  总结  /  未来规划", PURPLE),
]
for i, (num, title, desc, accent) in enumerate(items):
    y = Inches(3.0) + Inches(0.95) * i
    rect(s, Inches(1.0), y, Inches(0.06), Inches(0.75), accent)
    txt(s, Inches(1.3), y, Inches(1), Inches(0.7), num, size=40, color=accent, bold=True)
    txt(s, Inches(2.5), y + Inches(0.05), Inches(5), Inches(0.4), title, size=22, color=WHITE, bold=True)
    txt(s, Inches(2.5), y + Inches(0.45), Inches(8), Inches(0.3), desc, size=12, color=DIM)
rect(s, Inches(9.0), Inches(3.0), Inches(3.8), Inches(3.5), SURFACE)
txt(s, Inches(9.3), Inches(3.2), Inches(3.2), Inches(0.4), "📹 演示视频", size=13, color=CYAN, bold=True)
txt(s, Inches(9.3), Inches(3.8), Inches(3.2), Inches(2.5), "视频演示（AAA.md）：\n\n• 编辑模式绘画教学\n• 比赛模式实时标注\n• AI 快速/深度分析\n• Agent 建议完整闭环\n\nPPT 侧重：\nAgent 架构 + 商业卖点", size=10, color=DIM)
pg(s, 1)

# ============================
# 2. Part A 标题页
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s, Inches(0), Inches(0), Inches(0.25), H, CYAN)
txt(s, Inches(1.0), Inches(2.8), Inches(4), Inches(0.6), "PART  A", size=20, color=CYAN, bold=True)
txt(s, Inches(1.0), Inches(3.4), Inches(10), Inches(1), "为什么要做", size=52, color=WHITE, bold=True)
txt(s, Inches(1.0), Inches(4.6), Inches(10), Inches(0.5), "项目背景  ·  问题发现  ·  用户需求", size=16, color=DIM)
pg(s, 2)

# ============================
# 3. 问题 & 洞察
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.6), Inches(0.5), Inches(0.5))
txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4), "01  |  项目背景与问题发现", size=26, color=WHITE, bold=True)
rect(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.03), ORANGE)
txt(s, Inches(0.6), Inches(1.7), Inches(5.8), Inches(0.5), "💀  核心问题", size=20, color=ORANGE, bold=True)
bullets(s, Inches(0.6), Inches(2.3), Inches(5.8), Inches(2.5), [
    "CS2 死亡视角有 8-10 秒空闲——信息浪费",
    "业余战队无教练、无系统数据分析能力",
    "对手站位/投掷物/路线仅靠记忆不可靠",
    "现有工具只能画图，不具备战术理解能力",
], size=14)
rect(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.03), YELLOW)
txt(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(0.5), "💡  核心洞察", size=20, color=YELLOW, bold=True)
bullets(s, Inches(6.9), Inches(2.3), Inches(5.8), Inches(2.5), [
    "业余 IGL 需要能「读懂 CS2」的 AI",
    "通用 LLM 不懂地图点位和战术逻辑",
    "需要一个预调制的 CS2 专用 Agent",
    "用户自接 API Key → 低边际成本",
], size=14)
rect(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.8), SURFACE)
txt(s, Inches(1.0), Inches(5.35), Inches(11.3), Inches(0.5), '🎯 「Personal CS Coach 不是又一个战术画板——它是一个真正理解 CS2 地图、站位和战术的 AI Agent。」', size=17, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
footer(s); pg(s, 3)

# ============================
# 4. 用户 + 旅程
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.6), Inches(0.5), Inches(0.5))
txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4), "01  |  用户需求分析", size=26, color=WHITE, bold=True)
for i, (title, desc, icon) in enumerate([
    ("IGL（队内指挥）", "业余战队战术核心\n需要看懂对手习惯\n愿意为 AI 工具付费", "🎯"),
    ("进阶单排玩家", "高分段需要理解对手\n偏好低成本、低延迟方案\n自备 API Key 即可使用", "👤"),
    ("半职业 / 青训队", "有比赛需求预算有限\n需要复盘 + 数据积累\n重视 Agent 的多地图覆盖", "🏆"),
]):
    x = Inches(0.6) + Inches(4.1) * i
    rect(s, x, Inches(1.5), Inches(3.8), Inches(2.8), SURFACE)
    txt(s, x+Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.5), f"{icon}  {title}", size=15, color=CYAN, bold=True)
    txt(s, x+Inches(0.2), Inches(2.4), Inches(3.4), Inches(1.8), desc, size=12, color=DIM)
txt(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.4), "🗺  用户旅程", size=18, color=CYAN, bold=True)
journey = [
    ("① 赛前", "Agent 推荐\n预调制战术"),
    ("② 对局中", "第二屏幕\n快速标注"),
    ("③ 回合间隙", "Agent 分析\n对手倾向"),
    ("④ 暂停/中场", "深度分析\n生成对策"),
    ("⑤ 赛后", "Agent 学习\n优化战术库"),
]
for i, (step, desc) in enumerate(journey):
    x = Inches(0.6) + Inches(2.45) * i
    rect(s, x, Inches(5.2), Inches(2.2), Inches(1.6), SURFACE)
    txt(s, x+Inches(0.1), Inches(5.35), Inches(2.0), Inches(0.4), step, size=13, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.1), Inches(5.8), Inches(2.0), Inches(0.8), desc, size=11, color=DIM, align=PP_ALIGN.CENTER)
    if i < 4: txt(s, x+Inches(2.2), Inches(5.8), Inches(0.25), Inches(0.3), "→", size=16, color=CYAN, bold=True)
footer(s); pg(s, 4)

# ============================
# 5. Part B 标题页
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s, Inches(0), Inches(0), Inches(0.25), H, ORANGE)
txt(s, Inches(1.0), Inches(2.8), Inches(4), Inches(0.6), "PART  B", size=20, color=ORANGE, bold=True)
txt(s, Inches(1.0), Inches(3.4), Inches(10), Inches(1), "我们做了什么", size=52, color=WHITE, bold=True)
txt(s, Inches(1.0), Inches(4.6), Inches(10), Inches(0.5), "Agent 架构  ·  产品方案  ·  Demo 演示", size=16, color=DIM)
pg(s, 5)

# ============================
# 6. ⭐ 核心卖点 — Agent 架构（独立一页）
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.6), Inches(0.5), Inches(0.5), PURPLE)
txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4), "02  |  核心卖点：预调制 CS2 战术 Agent", size=26, color=WHITE, bold=True)

# Agent 概念图
rect(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.6), SURFACE)
txt(s, Inches(1.0), Inches(1.6), Inches(11), Inches(0.5), "🤖  Personal CS Coach Agent — 专为 CS2 预调制的战术智能体", size=20, color=PURPLE, bold=True)

agent_desc = """与通用 LLM 不同，Personal CS Coach Agent 在发货前已完成 CS2 领域预调制：

  ✅ 预训练  7 张竞赛地图的完整拓扑结构（点位/通路/掩体/投掷物弹道）
  ✅ 预注入  标准报点术语体系（A小/B小/拱门/长廊/水下/甜甜圈...）
  ✅ 预学习  经典战术模式（默认站位/Rush/慢攻/假打/双架/反清）
  ✅ 预校准  5 维评估模型（控图覆盖率/投掷物效率/影响力重叠/时间窗口/风险等级）

  💡 用户购买 Agent 后 → 自行接入 DeepSeek API → Agent 将 CS2 战术理解注入 LLM 推理"""
txt(s, Inches(1.0), Inches(2.2), Inches(11), Inches(1.6), agent_desc, size=12, color=DIM)

# 对比表
rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.8), SURFACE)
txt(s, Inches(1.0), Inches(4.35), Inches(11), Inches(0.4), "⚡  与通用 LLM 的差异", size=16, color=CYAN, bold=True)

# 表头
for i, (text, x) in enumerate([("", Inches(1.0)), ("通用 LLM", Inches(3.8)), ("Personal CS Coach Agent", Inches(7.8))]):
    txt(s, x, Inches(4.85), Inches(3), Inches(0.35), text, size=13, color=CYAN if i > 0 else WHITE, bold=True)

rows = [
    ("地图理解", "需在 Prompt 中描述:\n\"地图有 A/B 两个点...\"", "内置 7 张地图完整\n拓扑模型,直接推理"),
    ("报点术语", "不知道\"A小\"\"拱门\"\n\"香蕉道\"的位置", "预注入标准报点术语\n理解所有社区简称"),
    ("战术评估", "给通用评分,无\nCS2 领域知识校准", "5 维 CS2 专属模型\n按地图自适应调整"),
    ("成本", "每次推理 Token 消耗\n包含地图描述", "地图知识已内置\n推理 Token 更省"),
]
for i, (label, llm, agent) in enumerate(rows):
    y = Inches(5.25) + Inches(0.4) * i
    txt(s, Inches(1.0), y, Inches(2.5), Inches(0.35), label, size=11, color=WHITE)
    txt(s, Inches(3.8), y, Inches(3.5), Inches(0.35), llm, size=10, color=DIM)
    txt(s, Inches(7.8), y, Inches(4.5), Inches(0.35), agent, size=10, color=CYAN, bold=True)

footer(s); pg(s, 6)

# ============================
# 7. 产品方案（四大功能卡片）
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.6), Inches(0.5), Inches(0.5), ORANGE)
txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4), "02  |  产品功能矩阵", size=26, color=WHITE, bold=True)
cards = [
    ("🎨 画板标注", "7 种画笔工具\n箭头 / 圆圈 / 墙壁\n烟雾弹 / 闪光弹\n燃烧弹 / 手雷\n\nCT 蓝色 vs T 黄色\n拖拽即画 → 拖动修改"),
    ("⚡ 快速分析", "回合间隙 3-5s\nAgent 分析对手倾向\n\n置信度标注 30-85%\n一句话战术建议\n\n本地引擎 + 可选\nDeepSeek API 增强"),
    ("🔬 Agent 深度分析", "暂停 / 中场触发\n对手画像 + 柱状图\n3+ 条战术建议\n每条附带可采纳标记\n\n预训练模型确保\n建议符合 CS2 逻辑"),
    ("🤖 Agent 评估", "5 维雷达图\n（CS2 专属校准）\n控图 / 投掷 / 重叠\n时间窗口 / 风险\n\n优化前后对比\n采纳 Agent 方案"),
]
for i, (title, desc) in enumerate(cards):
    x = Inches(0.5) + Inches(3.15) * i
    rect(s, x, Inches(1.4), Inches(2.95), Inches(5.3), SURFACE)
    rect(s, x, Inches(1.4), Inches(2.95), Inches(0.03), CYAN)
    txt(s, x+Inches(0.2), Inches(1.6), Inches(2.5), Inches(0.5), title, size=15, color=CYAN, bold=True)
    txt(s, x+Inches(0.2), Inches(2.2), Inches(2.5), Inches(4.3), desc, size=11, color=DIM)
rect(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4), SURFACE)
txt(s, Inches(0.8), Inches(7.0), Inches(12), Inches(0.3), "技术栈：HTML + Canvas（纯前端）  |  Agent 预调制模型 + 用户自接 DeepSeek API  |  战术持久化：localStorage", size=10, color=DIM, align=PP_ALIGN.CENTER)
pg(s, 7)

# ============================
# 8. 原型展示（三大截图区）
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.6), Inches(0.5), Inches(0.5), ORANGE)
txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4), "02  |  产品原型展示", size=26, color=WHITE, bold=True)
for i, (title, desc, note) in enumerate([
    ("🏠 首页 — 地图选择", "7 张 CS2 竞赛地图\nMirage / Inferno / Dust2\nNuke / Ancient / Anubis\nOverpass", "📹 视频 0:30"),
    ("📐 战术编辑器", "拖拽绘图 + 标记拖动\n7 种画笔（含道具）\nAgent 评估雷达图\n优化对比弹窗", "📹 视频 0:30-2:00"),
    ("🎮 比赛仪表盘", "第二屏幕实时标注\nCT/T 颜色区分\nAgent 分析弹窗\n一键采纳 Agent 建议", "📹 视频 2:00-5:30"),
]):
    x = Inches(0.35) + Inches(4.25) * i
    rect(s, x, Inches(1.4), Inches(4.0), Inches(5.5), SURFACE)
    txt(s, x+Inches(0.2), Inches(1.55), Inches(3.6), Inches(0.5), title, size=14, color=CYAN, bold=True)
    txt(s, x+Inches(0.2), Inches(2.1), Inches(3.6), Inches(2.8), desc, size=11, color=DIM)
    rect(s, x+Inches(0.15), Inches(6.6), Inches(3.7), Inches(0.3), CYAN)
    txt(s, x+Inches(0.3), Inches(6.65), Inches(3.4), Inches(0.2), note, size=9, color=BG, align=PP_ALIGN.CENTER)
footer(s); pg(s, 8)

# ============================
# 9. Part C 标题页
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
rect(s, Inches(0), Inches(0), Inches(0.25), H, GREEN)
txt(s, Inches(1.0), Inches(2.8), Inches(4), Inches(0.6), "PART  C", size=20, color=GREEN, bold=True)
txt(s, Inches(1.0), Inches(3.4), Inches(10), Inches(1), "商业可行性", size=52, color=WHITE, bold=True)
txt(s, Inches(1.0), Inches(4.6), Inches(10), Inches(0.5), "Agent 卖点  ·  商业模式  ·  市场与竞品  ·  技术架构", size=16, color=DIM)
pg(s, 9)

# ============================
# 10. 商业模式 + Agent 卖点
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.6), Inches(0.5), Inches(0.5), GREEN)
txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4), "03  |  商业模式 & Agent 销售逻辑", size=26, color=WHITE, bold=True)

# 卖点大框
rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.8), SURFACE)
txt(s, Inches(0.9), Inches(1.45), Inches(11), Inches(0.4), "💎  产品形态", size=17, color=PURPLE, bold=True)
selling_points = [
    "用户购买的是一套 预调制的 CS2 战术 AI Agent — 不是 SaaS 订阅、不是原始 API。",
    "发货形态：前端画板工具（HTML+Canvas）+ Agent 配置文件（System Prompt / 7 地图 Prompt 模板 / 5 维评估权重）",
    "用户自行接入 DeepSeek API → Agent 将 CS2 战术知识注入 LLM 推理 → 本地规则引擎兜底",
]
bullets(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(1), selling_points, size=12, bullet='')

# 三级定价
tiers = [
    ("🆓  免费版", "$0", "单排玩家", ["基础 Agent（不含预调制）", "1 张地图", "本地规则引擎", "社区 Prompt 库"], DIM),
    ("⭐  Agent 版", "$9.9  一次性", "业余车队", ["完整预调制 Agent", "7 张地图全支持", "CS2 专属 5 维评估模型", "战术持久化 + 历史库"], CYAN),
    ("🏆  高阶版", "$29  年费", "半职业 / 青训", ["多队管理 Agent", "复盘数据导出", "API 优先分析 + 自定义", "Agent 持续更新"], YELLOW),
]
for i, (name, price, target, features, accent) in enumerate(tiers):
    x = Inches(0.4) + Inches(4.25) * i
    rect(s, x, Inches(3.3), Inches(4.0), Inches(3.8), SURFACE)
    txt(s, x+Inches(0.2), Inches(3.45), Inches(3.6), Inches(0.4), name, size=16, color=accent, bold=True)
    txt(s, x+Inches(0.2), Inches(3.8), Inches(3.6), Inches(0.6), price, size=24, color=accent, bold=True)
    txt(s, x+Inches(0.2), Inches(4.35), Inches(3.6), Inches(0.3), target, size=11, color=DIM)
    hr(s, x+Inches(0.2), Inches(4.7), Inches(3.6), accent, Inches(0.015))
    bullets(s, x+Inches(0.2), Inches(4.85), Inches(3.6), Inches(2), features, size=11, bullet='✓')

# 价值主张
rect(s, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.25), SURFACE)
txt(s, Inches(0.8), Inches(7.15), Inches(12), Inches(0.25), '💡 「一次购买，永久使用 — Agent 是你的战术教练，API Key 是你的调用通道」 — 一次性付费 + 自带 API Key', size=10, color=DIM, align=PP_ALIGN.CENTER)
pg(s, 10)

# ============================
# 11. 市场 + 竞品 + 架构
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.6), Inches(0.5), Inches(0.5), GREEN)
txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4), "03  |  市场分析 / 竞品对比 / 技术架构", size=26, color=WHITE, bold=True)

# 左—市场
rect(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.03), CYAN)
txt(s, Inches(0.5), Inches(1.55), Inches(6), Inches(0.4), "📊  市场规模", size=16, color=CYAN, bold=True)
bullets(s, Inches(0.5), Inches(2.1), Inches(6), Inches(2), [
    "CS2 日活约 80 万，月活约 2500 万",
    "业余战队 5-8 万支（社区联赛+校园赛）",
    "目标渗透率 5% = 2500-4000 车队",
    "Agent 版 TAM ≈ $40K（$9.9 × 4000）",
], size=12)

# 右—竞品
rect(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.03), CYAN)
txt(s, Inches(7.0), Inches(1.55), Inches(6), Inches(0.4), "⚔  竞品对比", size=16, color=CYAN, bold=True)
for i, (name, desc, status) in enumerate([
    ("CS2 Stratbook", "静态画板，无战术理解", "✗"),
    ("ChatGPT / DeepSeek", "通用 LLM，不认 CS2 点位", "✗"),
    ("Personal CS Coach Agent", "预调制 CS2 Agent + 可操作建议", "✓"),
]):
    y = Inches(2.15) + Inches(0.65) * i
    rect(s, Inches(7.0), y, Inches(6), Inches(0.55), SURFACE)
    txt(s, Inches(7.2), y+Inches(0.1), Inches(2.5), Inches(0.35), f"{status}  {name}", size=12, color=CYAN if status=='✓' else DIM, bold=status=='✓')
    txt(s, Inches(9.8), y+Inches(0.1), Inches(3), Inches(0.35), desc, size=11, color=DIM)

# 技术架构
rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8), SURFACE)
txt(s, Inches(0.8), Inches(4.35), Inches(12), Inches(0.4), "🔧  技术架构：Agent + API 双层", size=15, color=CYAN, bold=True)
arch = """┌──────────────────── 用户侧 ─────────────────────┐
│  HTML + Canvas 画板  ←→  localStorage 本地持久化  │
└──────────────────────┬────────────────────────────┘
                       │
  ┌────────────────────┴────────────────────────────┐
  │           Personal CS Coach Agent（预调制）              │
  │  ┌──────────┐  ┌──────────┐  ┌──────────────┐  │
  │  │7地图拓扑  │  │标准报点  │  │5维评估模型   │  │
  │  │预训练模型 │  │术语注入  │  │CS2专属校准   │  │
  │  └──────────┘  └──────────┘  └──────────────┘  │
  └──────────────────────┬────────────────────────────┘
                         │
           ┌─────────────┴─────────────┐
           │   用户自行接入 DeepSeek    │
           │   (API Key 本地存储)       │
           └───────────────────────────┘
                         │
           ┌─────────────┴─────────────┐
           │  离线兜底：本地规则引擎    │
           └───────────────────────────┘"""
txt(s, Inches(1.0), Inches(4.85), Inches(12), Inches(2), arch, size=10, color=DIM)
footer(s); pg(s, 11)

# ============================
# 12. 总结 + Agent 路线图
# ============================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
bar(s, Inches(0.6), Inches(0.5), Inches(0.5), PURPLE)
txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.4), "04  |  项目总结 & Agent 预训练路线图", size=26, color=WHITE, bold=True)

# 左—已完成
txt(s, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4), "✅  已完成（MVP）", size=18, color=CYAN, bold=True)
bullets(s, Inches(0.6), Inches(1.9), Inches(6), Inches(2.5), [
    "7 张地图完整画板 + 7 种画笔",
    "CT/T 颜色分离 + 拖拽绘图 + 标记拖动",
    "本地规则引擎（5 维评分 + 分析）",
    "API 接入层（OpenAI/DeepSeek/Claude）",
    "战术持久化（localStorage）",
    "Agent 建议可采纳画到画板",
], size=12)

# 右—Agent 预训练路线图
txt(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.4), "🗓  Agent 预训练路线图", size=18, color=PURPLE, bold=True)
roadmap = [
    ("Phase 1", "MVP 验证", "✅", "产品可用性 + Demo 闭环"),
    ("Phase 2", "地图 Prompt 注入", "⏳", "7 张地图点位/术语写入 System Prompt"),
    ("Phase 3", "战术模式库", "⬜", "Rush/慢攻/假打/反清等经典模式预学习"),
    ("Phase 4", "用户反馈调优", "⬜", "收集标注数据 → 校准 Agent 评估权重"),
    ("Phase 5", "发布 Agent v1.0", "⬜", "打包 Agent 配置 + 用户自接 API"),
]
for i, (phase, title, status, desc) in enumerate(roadmap):
    y = Inches(1.95) + Inches(0.65) * i
    rect(s, Inches(7.0), y, Inches(6), Inches(0.55), SURFACE)
    txt(s, Inches(7.15), y+Inches(0.1), Inches(0.7), Inches(0.35), f"{status}  {phase}", size=10, color=CYAN if status=='✅' else DIM, bold=True)
    txt(s, Inches(8.0), y+Inches(0.1), Inches(2), Inches(0.35), title, size=12, color=WHITE)
    txt(s, Inches(10.0), y+Inches(0.1), Inches(2.8), Inches(0.35), desc, size=10, color=DIM)

# 底部总结
rect(s, Inches(0.6), Inches(4.8), Inches(12.3), Inches(1.5), SURFACE)
txt(s, Inches(1.0), Inches(5.0), Inches(11.5), Inches(0.4), "💬  总结", size=18, color=CYAN, bold=True)
summary = (
    "Personal CS Coach 的核心壁垒不是画板——是预调制的 CS2 战术 Agent。\n\n"
    "当通用 LLM 还在问「A 小道和 B 小道在哪」时，我们的 Agent 已经知道 Mirage 的拱门到 A 点的最短路径、"
    "Inferno 的香蕉道是 B 点防守的生命线、Dust2 的 A 大门是 T 方进攻的必争之地。\n\n"
    "这是 CS2 战术分析从「工具」到「智能体」的范式转变。"
)
txt(s, Inches(1.0), Inches(5.45), Inches(11.5), Inches(0.8), summary, size=11, color=DIM)

# 致谢
rect(s, Inches(0.6), Inches(6.5), Inches(12.3), Inches(0.65), SURFACE)
txt(s, Inches(1.0), Inches(6.6), Inches(11.5), Inches(0.45), "🙏  感谢聆听  ·  欢迎提问", size=28, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(6.95), Inches(11.5), Inches(0.2), "📹 视频演示见 AAA.md  |  🎮 源码：HTML+Canvas 纯前端  |  🤖 卖点：预调制 Agent + 用户自接 DeepSeek", size=10, color=DIM, align=PP_ALIGN.CENTER)
pg(s, 12)

# ============================
# 保存
# ============================
path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'Personal_CS_Coach_答辩PPT.pptx')
prs.save(path)
print(f'✅ PPT 已生成: {path}')
print(f'   共 {len(prs.slides)} 页  |  16:9 宽屏')
print(f'   核心卖点：预调制 CS2 战术 Agent + 用户自接 DeepSeek API')